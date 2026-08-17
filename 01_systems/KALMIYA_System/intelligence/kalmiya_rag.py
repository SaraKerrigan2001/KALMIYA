"""
kalmiya_rag.py — RAG (Retrieval Augmented Generation) para KALMIYA
===================================================================
Permite a KALMIYA buscar en TUS documentos antes de responder:
  - Notas de Obsidian (.md)
  - PDFs del SENA / ADSO
  - Documentos .txt, .docx
  - Código fuente .py, .js, .java

Flujo RAG:
  1. Indexar documentos → dividir en chunks → guardar en ChromaDB local
  2. Al preguntar → buscar chunks relevantes → inyectar en el prompt
  3. KALMIYA responde con contexto real de TUS documentos

Sin API key externa — todo funciona localmente con sentence-transformers.
"""

import os, sys, json, hashlib, re
os.environ["HF_HUB_OFFLINE"] = "1"
from pathlib import Path
from datetime import datetime
from typing import Optional

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from database import log_command, update_memory, get_memory
from _logging import get_logger

logger = get_logger(__name__)

# ── Dependencias opcionales ────────────────────────────────────────────────────
try:
    import chromadb
    from chromadb.config import Settings
    CHROMA_OK = True
except ImportError:
    CHROMA_OK = False
    logger.warning("[RAG] chromadb no disponible — instala: pip install chromadb")

try:
    from sentence_transformers import SentenceTransformer
    ST_OK = True
except ImportError:
    ST_OK = False
    logger.warning("[RAG] sentence-transformers no disponible — instala: pip install sentence-transformers")

try:
    import pypdf
    PYPDF_OK = True
except ImportError:
    PYPDF_OK = False

try:
    from docx import Document
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

# ── Configuración ──────────────────────────────────────────────────────────────
RAG_DIR     = Path(__file__).parent / "rag_db"
VAULT_PATH  = Path(r"c:\Users\maria\env")
CHUNK_SIZE  = 500    # caracteres por chunk
CHUNK_OVERLAP = 50   # solapamiento entre chunks
TOP_K       = 5      # chunks más relevantes a recuperar
MIN_SCORE   = 0.3    # similitud mínima para incluir un chunk

# Extensiones de archivo indexables
EXTENSIONES_INDEXABLES = {
    ".md", ".txt", ".py", ".js", ".java", ".html",
    ".css", ".json", ".yaml", ".yml", ".pdf", ".docx", ".csv", ".xlsx", ".pptx"
}

# Carpetas a excluir del indexado
EXCLUIR_CARPETAS = {
    ".obsidian", "__pycache__", ".git", "node_modules",
    "02_infrastructure", "site-packages", "_BACKUPS",
    "rag_db", "temp_audio", ".pytest_cache"
}

# ── Estado global ──────────────────────────────────────────────────────────────
_cliente_chroma = None
_coleccion      = None
_modelo_embed   = None
_indexado       = False

# ══════════════════════════════════════════════════════════════════════════════
# INICIALIZACIÓN
# ══════════════════════════════════════════════════════════════════════════════

def _init_rag() -> bool:
    """Inicializa ChromaDB y el modelo de embeddings."""
    global _cliente_chroma, _coleccion, _modelo_embed

    if not CHROMA_OK:
        return False

    try:
        RAG_DIR.mkdir(parents=True, exist_ok=True)
        _cliente_chroma = chromadb.PersistentClient(
            path=str(RAG_DIR),
            settings=Settings(anonymized_telemetry=False)
        )
        _coleccion = _cliente_chroma.get_or_create_collection(
            name="kalmiya_docs",
            metadata={"hnsw:space": "cosine"}
        )
        logger.info(f"[RAG] ChromaDB iniciado — {_coleccion.count()} documentos indexados")
    except Exception as e:
        logger.error(f"[RAG] Error iniciando ChromaDB: {e}")
        return False

    if ST_OK:
        try:
            # Modelo multilingüe ligero — funciona en español
            _modelo_embed = SentenceTransformer(
                "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2",
                local_files_only=True
            )
            logger.info("[RAG] Modelo de embeddings cargado")
        except Exception as e:
            logger.warning(f"[RAG] No se pudo cargar sentence-transformers: {e}")
    else:
        logger.warning("[RAG] Usando embeddings simples (sin sentence-transformers)")

    return True


def _get_embedding(texto: str) -> list[float]:
    """Genera el embedding de un texto."""
    if _modelo_embed is not None:
        try:
            return _modelo_embed.encode(texto, normalize_embeddings=True).tolist()
        except Exception:
            pass
    # Fallback mejorado: TF-IDF simple con vocabulario fijo
    # Produce vectores de 384 dimensiones más informativos que un hash
    import math
    palabras = re.sub(r'[^a-záéíóúñü\s]', '', texto.lower()).split()
    vec = [0.0] * 384
    for palabra in palabras:
        # Usar hash de la palabra para distribuir en el vector
        h = int(hashlib.md5(palabra.encode()).hexdigest(), 16)
        idx = h % 384
        vec[idx] += 1.0
    # Normalizar
    norm = math.sqrt(sum(x*x for x in vec)) or 1.0
    return [x / norm for x in vec]


# ══════════════════════════════════════════════════════════════════════════════
# PROCESAMIENTO DE DOCUMENTOS
# ══════════════════════════════════════════════════════════════════════════════

def _leer_archivo(ruta: Path) -> str:
    """Lee el contenido de un archivo según su extensión de forma segura."""
    ext = ruta.suffix.lower()
    try:
        if ext == ".pdf":
            if not PYPDF_OK:
                logger.warning(f"[RAG] No se puede leer {ruta.name}: 'pypdf' no está disponible.")
                return ""
            texto_pdf = []
            with open(ruta, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    texto_pdf.append(page.extract_text() or "")
            return "\n".join(texto_pdf)
        elif ext == ".docx":
            if not DOCX_OK:
                logger.warning(f"[RAG] No se puede leer {ruta.name}: 'python-docx' no está disponible.")
                return ""
            doc = Document(str(ruta))
            return "\n".join(p.text for p in doc.paragraphs)
        else:
            # New handlers for CSV, Excel, PowerPoint
            if ext == ".csv":
                try:
                    import pandas as pd
                    df = pd.read_csv(ruta)
                    return df.to_csv(index=False)
                except Exception as e:
                    logger.warning(f"[RAG] No se puede leer CSV {ruta.name}: {e}")
                    return ""
            elif ext == ".xlsx":
                try:
                    import pandas as pd
                    df = pd.read_excel(ruta, engine='openpyxl')
                    return df.to_csv(index=False)
                except Exception as e:
                    logger.warning(f"[RAG] No se puede leer Excel {ruta.name}: {e}")
                    return ""
            elif ext == ".pptx":
                try:
                    # pyrefly: ignore [missing-import]
                    from pptx import Presentation
                    prs = Presentation(ruta)
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                texts.append(shape.text)
                    return "\n".join(texts)
                except Exception as e:
                    logger.warning(f"[RAG] No se puede leer PowerPoint {ruta.name}: {e}")
                    return ""
            else:
                return ruta.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        logger.error(f"[RAG] Error leyendo el archivo {ruta}: {e}")
        return ""


def _dividir_en_chunks(texto: str, fuente: str) -> list[dict]:
    """
    Divide un texto en chunks usando RecursiveCharacterTextSplitter de LangChain.
    Mejora enormemente la cohesión semántica de los fragmentos en comparación al método manual.
    """
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except ImportError:
        logger.warning("[RAG] LangChain no disponible. Usando fallback básico.")
        # Fallback básico si LangChain aún no termina de instalarse
        chunks_basicos = []
        for i in range(0, len(texto), CHUNK_SIZE):
            chunk = texto[i:i+CHUNK_SIZE]
            if len(chunk.strip()) > 50:
                chunks_basicos.append({
                    "id": hashlib.md5(f"{fuente}_{i}".encode()).hexdigest(),
                    "texto": chunk, "fuente": fuente, "indice": i, "longitud": len(chunk)
                })
        return chunks_basicos

    texto = re.sub(r'\n{3,}', '\n\n', texto.strip())
    
    # Usamos el separador recursivo que es el estándar de la industria
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", r"(?<=\. )", " ", ""]
    )
    
    textos_divididos = splitter.split_text(texto)
    
    chunks = []
    for idx, chunk_texto in enumerate(textos_divididos):
        if len(chunk_texto.strip()) > 50:
            chunk_id = hashlib.md5(f"{fuente}_{idx}".encode()).hexdigest()
            chunks.append({
                "id": chunk_id,
                "texto": chunk_texto,
                "fuente": fuente,
                "indice": idx,
                "longitud": len(chunk_texto),
            })
            
    return chunks


def _subdividir_chunk(texto: str, fuente: str, start_idx: int) -> list[dict]:
    """Deprecado: RecursiveCharacterTextSplitter ya maneja el tamaño máximo."""
    return []


def indexar_documento(ruta: Path, forzar: bool = False) -> int:
    """
    Indexa un documento en ChromaDB.

    Args:
        ruta:    Ruta al archivo.
        forzar:  Reindexar aunque ya exista.

    Returns:
        Número de chunks indexados.
    """
    if not _coleccion:
        if not _init_rag():
            return 0

    fuente = str(ruta.relative_to(VAULT_PATH)) if ruta.is_relative_to(VAULT_PATH) else ruta.name

    # Verificar si ya está indexado (por hash del archivo)
    try:
        contenido = _leer_archivo(ruta)
        if not contenido.strip():
            return 0

        file_hash = hashlib.md5(contenido.encode()).hexdigest()
        # Sanitizar clave: reemplazar /, . y otros caracteres inválidos
        sanitized = fuente.replace('/', '_').replace('.', '_').replace(' ', '_')
        hash_key  = f"rag_hash_{sanitized}"
        guardado  = get_memory(hash_key)

        if guardado == file_hash and not forzar:
            return 0  # Sin cambios

        # Eliminar chunks anteriores del mismo archivo
        try:
            existing = _coleccion.get(where={"fuente": fuente})
            if existing["ids"]:
                _coleccion.delete(ids=existing["ids"])
        except Exception:
            pass

        # Generar nuevos chunks
        chunks = _dividir_en_chunks(contenido, fuente)
        if not chunks:
            return 0

        # Generar embeddings e insertar
        ids        = [c["id"] for c in chunks]
        textos     = [c["texto"] for c in chunks]
        embeddings = [_get_embedding(t) for t in textos]
        metadatos  = [{
            "fuente":  c["fuente"],
            "indice":  c["indice"],
            "tipo":    ruta.suffix.lower(),
            "nombre":  ruta.name,
            "fecha":   datetime.fromtimestamp(ruta.stat().st_mtime).isoformat(),
        } for c in chunks]

        _coleccion.add(
            ids=ids,
            embeddings=embeddings,
            documents=textos,
            metadatas=metadatos,
        )

        update_memory(hash_key, file_hash)
        logger.debug(f"[RAG] Indexado: {fuente} — {len(chunks)} chunks")
        return len(chunks)

    except Exception as e:
        logger.warning(f"[RAG] Error indexando {ruta}: {e}")
        return 0


def indexar_vault(carpeta: Path = None, mostrar_progreso: bool = True) -> dict:
    """
    Indexa todos los documentos del vault de Obsidian y carpeta KALMIYA.

    Returns:
        dict con estadísticas del indexado.
    """
    if not _init_rag():
        return {"error": "RAG no disponible"}

    carpeta = carpeta or VAULT_PATH
    stats   = {"archivos": 0, "chunks": 0, "omitidos": 0, "errores": 0}

    archivos = [
        p for p in carpeta.rglob("*")
        if p.is_file()
        and p.suffix.lower() in EXTENSIONES_INDEXABLES
        and not any(ex in p.parts for ex in EXCLUIR_CARPETAS)
    ]

    if mostrar_progreso:
        print(f"\n  [RAG] Indexando {len(archivos)} archivos...")

    for i, archivo in enumerate(archivos):
        try:
            n = indexar_documento(archivo)
            if n > 0:
                stats["archivos"] += 1
                stats["chunks"]   += n
                if mostrar_progreso and i % 20 == 0:
                    print(f"  [RAG] Progreso: {i+1}/{len(archivos)} — {stats['chunks']} chunks")
            else:
                stats["omitidos"] += 1
        except Exception:
            stats["errores"] += 1

    update_memory("rag_ultimo_indexado", datetime.now().isoformat())
    update_memory("rag_total_chunks",    str(stats["chunks"]))

    if mostrar_progreso:
        print(f"\n  [RAG] ✅ Indexado completo:")
        print(f"        {stats['archivos']} archivos nuevos/modificados")
        print(f"        {stats['chunks']} chunks en total")
        print(f"        {_coleccion.count()} chunks en la base vectorial")

    log_command("[RAG] Indexado", json.dumps(stats), source="system")
    return stats

# ══════════════════════════════════════════════════════════════════════════════
# BÚSQUEDA RAG CON MULTI-QUERY Y RE-RANKING
# ══════════════════════════════════════════════════════════════════════════════

def _expandir_query(query: str) -> list[str]:
    """Genera variaciones de la query para mejorar el recall."""
    queries = [query]
    # Variación 1: sin signos de interrogación ni tildes sueltas
    limpia = re.sub(r'[¿?¡!]', '', query).strip()
    if limpia != query:
        queries.append(limpia)
    # Variación 2: palabras clave principales (>3 chars)
    palabras = [p for p in re.sub(r'[^\w\s]', '', query.lower()).split() if len(p) > 3]
    if len(palabras) >= 2:
        queries.append(' '.join(palabras))
    return queries[:3]  # Máximo 3 variaciones


def _rerank_chunks(chunks: list[dict], query: str) -> list[dict]:
    """
    Re-rankea chunks considerando:
    1. Similitud coseno (ya calculada)
    2. Boost si la query aparece literalmente en el texto
    3. Boost por recencia del documento
    """
    query_lower = query.lower()
    palabras_query = set(re.sub(r'[^\w\s]', '', query_lower).split())

    for chunk in chunks:
        score = chunk["similitud"]

        # Boost por keywords: +0.05 por cada palabra de la query que aparece
        texto_lower = chunk["texto"].lower()
        keywords_encontradas = sum(1 for p in palabras_query if p in texto_lower and len(p) > 3)
        score += keywords_encontradas * 0.05

        # Boost por match exacto de la query
        if query_lower in texto_lower:
            score += 0.15

        # Boost por recencia (documentos recientes tienen más peso)
        try:
            fecha_doc = datetime.fromisoformat(chunk.get("fecha", ""))
            dias_antiguedad = (datetime.now() - fecha_doc).days
            if dias_antiguedad < 7:
                score += 0.08
            elif dias_antiguedad < 30:
                score += 0.04
        except (ValueError, TypeError):
            pass

        chunk["score_final"] = round(min(score, 1.0), 3)

    chunks.sort(key=lambda x: x.get("score_final", x["similitud"]), reverse=True)
    return chunks


def buscar_rag(query: str, top_k: int = TOP_K,
               filtro_tipo: str = None) -> list[dict]:
    """
    Busca los chunks más relevantes usando multi-query + re-ranking.

    Args:
        query:       Pregunta o consulta del usuario.
        top_k:       Número de resultados a devolver.
        filtro_tipo: Filtrar por tipo de archivo (.md, .py, .pdf...)

    Returns:
        Lista de chunks relevantes ordenados por score final.
    """
    if not _coleccion:
        if not _init_rag():
            return []

    if _coleccion.count() == 0:
        return []

    try:
        # Multi-query: expandir la consulta
        queries = _expandir_query(query)
        todos_chunks = {}
        where = {"tipo": filtro_tipo} if filtro_tipo else None
        fetch_k = min(top_k * 2, _coleccion.count())  # Recuperar más para re-rankear

        for q in queries:
            embedding = _get_embedding(q)
            resultados = _coleccion.query(
                query_embeddings=[embedding],
                n_results=fetch_k,
                where=where,
                include=["documents", "metadatas", "distances"]
            )

            for i, doc in enumerate(resultados["documents"][0]):
                distancia = resultados["distances"][0][i]
                similitud = 1 - distancia
                metadato = resultados["metadatas"][0][i]
                chunk_key = metadato.get("fuente", "") + "_" + str(metadato.get("indice", i))

                if similitud >= MIN_SCORE and chunk_key not in todos_chunks:
                    todos_chunks[chunk_key] = {
                        "texto":    doc,
                        "fuente":   metadato.get("fuente", "?"),
                        "nombre":   metadato.get("nombre", "?"),
                        "tipo":     metadato.get("tipo", "?"),
                        "similitud": round(similitud, 3),
                        "fecha":    metadato.get("fecha", ""),
                    }

        chunks = list(todos_chunks.values())

        # Re-ranking con boost por keywords y recencia
        chunks = _rerank_chunks(chunks, query)

        # Devolver solo top_k
        resultado = chunks[:top_k]
        logger.debug(f"[RAG] '{query[:40]}...' → {len(resultado)} chunks (de {len(chunks)} candidatos)")
        return resultado

    except Exception as e:
        logger.warning(f"[RAG] Error en búsqueda: {e}")
        return []


def buscar_rag_manual(query: str, top_k: int = 10) -> list[dict]:
    """Búsqueda manual para la UI — devuelve más resultados con preview."""
    return buscar_rag(query, top_k=top_k)


def construir_contexto_rag(query: str, top_k: int = TOP_K) -> tuple:
    """
    Construye el bloque de contexto RAG para inyectar en el prompt.

    Returns:
        Tuple de (contexto_str, lista_fuentes) — contexto formateado y fuentes usadas.
    """
    chunks = buscar_rag(query, top_k=top_k)
    if not chunks:
        return "", []

    lineas = [
        "CONTEXTO DE TUS DOCUMENTOS PERSONALES:",
        "(Usa esta información para responder. Cita la fuente cuando sea relevante.)"
    ]
    fuentes_vistas = set()
    fuentes_lista = []

    for chunk in chunks:
        fuente = chunk["fuente"]
        score = chunk.get("score_final", chunk["similitud"])
        if fuente not in fuentes_vistas:
            fuentes_vistas.add(fuente)
            fuentes_lista.append({"fuente": fuente, "score": score, "tipo": chunk["tipo"]})
            lineas.append(f"\n📄 Fuente: {fuente} (relevancia: {score})")
        lineas.append(chunk["texto"])

    lineas.append("\n--- Fin del contexto ---")
    lineas.append("Responde basándote en este contexto. Si citas información específica, menciona de qué documento proviene.")

    return "\n".join(lineas), fuentes_lista


def responder_con_rag(query: str, force_engine: str = "") -> dict:
    """
    Responde una pregunta usando RAG + IA:
    1. Busca contexto relevante en los documentos
    2. Inyecta ese contexto en el prompt
    3. Pide a la IA que responda usando ese contexto

    Returns:
        Dict con 'response' (texto), 'rag_sources' (lista fuentes), 'rag_used' (bool).
    """
    from brain import ask_kalmiya, _conversation_history

    contexto, fuentes = construir_contexto_rag(query, top_k=TOP_K)
    rag_used = False
    if contexto:
        prompt = f"{contexto}\n\nPregunta del usuario: {query}"
        logger.info(f"[RAG] Respondiendo con contexto ({len(contexto)} chars, {len(fuentes)} fuentes)")
        rag_used = True
    else:
        prompt = query
        logger.debug("[RAG] Sin contexto relevante — respuesta directa")

    respuesta = ask_kalmiya(prompt, force_engine=force_engine)
    return {
        "response": respuesta,
        "rag_sources": fuentes,
        "rag_used": rag_used,
    }


def get_rag_stats() -> dict:
    """Devuelve estadísticas del sistema RAG con distribución por tipo."""
    if not _coleccion:
        _init_rag()

    total_chunks = _coleccion.count() if _coleccion else 0
    ultimo_idx = get_memory("rag_ultimo_indexado") or "Nunca"

    # Distribución por tipo de archivo
    tipo_stats = {}
    if _coleccion and total_chunks > 0:
        try:
            todos = _coleccion.get(include=["metadatas"])
            for meta in todos.get("metadatas", []):
                tipo = meta.get("tipo", "otro")
                tipo_stats[tipo] = tipo_stats.get(tipo, 0) + 1
        except Exception:
            pass

    return {
        "disponible":       CHROMA_OK and ST_OK,
        "chroma_ok":        CHROMA_OK,
        "embeddings_ok":    ST_OK,
        "chunks_en_db":     total_chunks,
        "ultimo_indexado":  ultimo_idx,
        "vault_path":       str(VAULT_PATH),
        "extensiones":      sorted(EXTENSIONES_INDEXABLES),
        "por_tipo":         tipo_stats,
    }


def limpiar_base_vectorial() -> dict:
    """Limpia completamente la base vectorial de ChromaDB."""
    global _coleccion
    if not _coleccion:
        if not _init_rag():
            return {"error": "RAG no disponible"}
    try:
        count_antes = _coleccion.count()
        # Borrar y recrear la colección
        _cliente_chroma.delete_collection("kalmiya_docs")
        _coleccion = _cliente_chroma.get_or_create_collection(
            name="kalmiya_docs",
            metadata={"hnsw:space": "cosine"}
        )
        logger.info(f"[RAG] Base vectorial limpiada — {count_antes} chunks eliminados")
        return {"resultado": f"Base limpiada: {count_antes} chunks eliminados", "chunks_eliminados": count_antes}
    except Exception as e:
        logger.error(f"[RAG] Error limpiando base vectorial: {e}")
        return {"error": str(e)}


def indexar_url(url: str) -> dict:
    """
    Descarga una página web, extrae texto limpio y lo indexa en ChromaDB.

    Args:
        url: URL de la página a indexar.

    Returns:
        dict con estadísticas del indexado.
    """
    if not _coleccion:
        if not _init_rag():
            return {"error": "RAG no disponible"}

    try:
        import urllib.request
        from html.parser import HTMLParser

        # Descargar HTML
        req = urllib.request.Request(url, headers={'User-Agent': 'KALMIYA-RAG/1.0'})
        with urllib.request.urlopen(req, timeout=15) as resp:
            html = resp.read().decode('utf-8', errors='ignore')

        # Extraer texto limpio del HTML
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.textos = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                self._skip = tag in ('script', 'style', 'nav', 'footer', 'header')

            def handle_endtag(self, tag):
                if tag in ('script', 'style', 'nav', 'footer', 'header'):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip:
                    texto = data.strip()
                    if texto:
                        self.textos.append(texto)

        extractor = TextExtractor()
        extractor.feed(html)
        contenido = '\n'.join(extractor.textos)

        if not contenido.strip() or len(contenido) < 50:
            return {"error": "No se pudo extraer texto de la URL"}

        # Generar chunks
        fuente = f"web:{url[:80]}"
        chunks = _dividir_en_chunks(contenido, fuente)
        if not chunks:
            return {"error": "Contenido muy corto para indexar"}

        # Eliminar indexado anterior de la misma URL
        try:
            existing = _coleccion.get(where={"fuente": fuente})
            if existing["ids"]:
                _coleccion.delete(ids=existing["ids"])
        except Exception:
            pass

        # Indexar
        ids = [c["id"] for c in chunks]
        textos = [c["texto"] for c in chunks]
        embeddings = [_get_embedding(t) for t in textos]
        metadatos = [{
            "fuente": fuente,
            "indice": c["indice"],
            "tipo": ".web",
            "nombre": url[:100],
            "fecha": datetime.now().isoformat(),
        } for c in chunks]

        _coleccion.add(ids=ids, embeddings=embeddings, documents=textos, metadatas=metadatos)

        logger.info(f"[RAG] URL indexada: {url[:60]} — {len(chunks)} chunks")
        return {"url": url, "chunks": len(chunks), "caracteres": len(contenido)}

    except Exception as e:
        logger.error(f"[RAG] Error indexando URL {url}: {e}")
        return {"error": str(e)}


def imprimir_rag_stats():
    """Imprime estadísticas del RAG en consola."""
    s = get_rag_stats()
    print("\n╔" + "═"*55 + "╗")
    print("║" + "  📚  KALMIYA RAG — Estado del Sistema".center(55) + "║")
    print("╠" + "═"*55 + "╣")
    print(f"║  Disponible      : {'✅ Sí' if s['disponible'] else '❌ No':<39}║")
    print(f"║  ChromaDB        : {'✅' if s['chroma_ok'] else '❌ pip install chromadb':<39}║")
    print(f"║  Embeddings      : {'✅ sentence-transformers' if s['embeddings_ok'] else '⚠️  modo básico':<39}║")
    print(f"║  Chunks en DB    : {str(s['chunks_en_db']):<39}║")
    print(f"║  Último indexado : {str(s['ultimo_indexado'])[:39]:<39}║")
    print(f"║  Vault           : {str(s['vault_path'])[:39]:<39}║")
    if s.get('por_tipo'):
        print("╠" + "═"*55 + "╣")
        for tipo, count in sorted(s['por_tipo'].items(), key=lambda x: -x[1]):
            print(f"║  {tipo:<16}: {str(count):<37}║")
    print("╚" + "═"*55 + "╝\n")
